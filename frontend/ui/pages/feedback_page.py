from __future__ import annotations

from pathlib import Path
import re
import zipfile
from xml.etree import ElementTree

from PySide6.QtCore import Qt
from PySide6.QtWidgets import QFileDialog, QHBoxLayout, QLabel, QListWidget, QListWidgetItem, QPlainTextEdit, QVBoxLayout, QWidget

try:
	from pypdf import PdfReader
except ImportError:  # pragma: no cover
	PdfReader = None  # type: ignore[assignment]

try:
	from docx import Document
except ImportError:  # pragma: no cover
	Document = None  # type: ignore[assignment]

try:
	from pptx import Presentation
except ImportError:  # pragma: no cover
	Presentation = None  # type: ignore[assignment]

try:
	import olefile
except ImportError:  # pragma: no cover
	olefile = None  # type: ignore[assignment]

from ...components.buttons import AppButton
from ...components.cards import CardWidget
from ...api_common import ApiError
from ...controllers import AgentController


class FeedbackPage(QWidget):
	def __init__(self, parent: QWidget | None = None) -> None:
		super().__init__(parent)
		self._uploaded_files: list[Path] = []
		self._file_ids: dict[str, str] = {}
		self._controller = AgentController()

		root = QVBoxLayout(self)
		root.setContentsMargins(0, 0, 0, 0)
		root.setSpacing(12)

		summary = CardWidget("피드백")
		summary_text = QLabel("문서를 업로드하면 파일별 피드백을 자동으로 생성합니다.")
		summary_text.setObjectName("PageSubtitle")
		summary_text.setWordWrap(True)

		stats = QLabel("업로드 파일: 0    |    분석 상태: 대기")
		stats.setObjectName("WarningSummary")
		self._stats = stats

		button_row = QHBoxLayout()
		button_row.setSpacing(8)

		upload_btn = AppButton("문서 업로드")
		upload_btn.clicked.connect(self._upload_documents)

		clear_btn = AppButton("목록 비우기", variant="ghost")
		clear_btn.clicked.connect(self._clear_documents)

		button_row.addWidget(upload_btn)
		button_row.addWidget(clear_btn)
		button_row.addStretch(1)

		summary.layout.addWidget(summary_text)
		summary.layout.addWidget(stats)
		summary.layout.addLayout(button_row)
		root.addWidget(summary)

		file_card = CardWidget("업로드 문서")
		file_hint = QLabel("파일을 선택하면 자동으로 피드백 결과가 생성됩니다.")
		file_hint.setObjectName("PageSubtitle")
		file_hint.setWordWrap(True)

		self.file_list = QListWidget()
		self.file_list.setObjectName("FeedbackFileList")
		self.file_list.setMinimumHeight(180)
		self.file_list.currentRowChanged.connect(self._on_file_selected)

		file_card.layout.addWidget(file_hint)
		file_card.layout.addWidget(self.file_list)
		root.addWidget(file_card)

		feedback_card = CardWidget("파일 피드백")
		feedback_hint = QLabel("선택한 파일의 위험 요소와 개선 제안을 확인하세요.")
		feedback_hint.setObjectName("PageSubtitle")
		feedback_hint.setWordWrap(True)

		self.feedback_output = QPlainTextEdit()
		self.feedback_output.setReadOnly(True)
		self.feedback_output.setObjectName("FeedbackOutput")
		self.feedback_output.setMinimumHeight(260)
		self.feedback_output.setPlainText("아직 업로드된 문서가 없습니다.")

		feedback_card.layout.addWidget(feedback_hint)
		feedback_card.layout.addWidget(self.feedback_output)
		root.addWidget(feedback_card)

		root.addStretch(1)

	def _upload_documents(self) -> None:
		files, _ = QFileDialog.getOpenFileNames(
			self,
			"피드백 문서 업로드",
			"",
			"Documents (*.txt *.md *.rst *.log *.pdf *.docx *.ppt *.pptx *.hwp);;All Files (*.*)",
		)
		if not files:
			return

		existing = {str(path) for path in self._uploaded_files}
		for file_path in files:
			if file_path in existing:
				continue
			path = Path(file_path)
			self._uploaded_files.append(path)
			existing.add(file_path)
			item = QListWidgetItem(path.name)
			item.setData(Qt.UserRole, str(path))
			self.file_list.addItem(item)

		if self.file_list.count() > 0 and self.file_list.currentRow() < 0:
			self.file_list.setCurrentRow(0)

		self._analyze_uploaded_documents([Path(file_path) for file_path in files])
		self._refresh_stats()

	def _clear_documents(self) -> None:
		self._uploaded_files.clear()
		self._file_ids.clear()
		self.file_list.clear()
		self.feedback_output.setPlainText("아직 업로드된 문서가 없습니다.")
		self._refresh_stats()

	def _on_file_selected(self, row: int) -> None:
		if row < 0 or row >= len(self._uploaded_files):
			return

		file_path = self._uploaded_files[row]
		file_id = self._file_ids.get(str(file_path))
		if not file_id:
			self.feedback_output.setPlainText(self._build_feedback(file_path))
			return

		try:
			result = self._controller.get_feedback_result(file_id)
			self.feedback_output.setPlainText(self._format_agent_feedback(result))
		except ApiError as e:
			self.feedback_output.setPlainText(f"API 요청 실패: {e}")

	def _analyze_uploaded_documents(self, paths: list[Path]) -> None:
		new_paths = [path for path in paths if str(path) not in self._file_ids]
		if not new_paths:
			return

		self.feedback_output.setPlainText("backend agent가 문서를 분석하는 중입니다...")
		try:
			uploaded = self._controller.upload_feedback_files(new_paths)
			file_ids: list[str] = []
			for path, item in zip(new_paths, uploaded):
				file_id = str(item.get("fileId") or "")
				if file_id:
					self._file_ids[str(path)] = file_id
					file_ids.append(file_id)
			if file_ids:
				self._controller.analyze_feedback(file_ids)
			if self.file_list.currentRow() >= 0:
				self._on_file_selected(self.file_list.currentRow())
		except ApiError as e:
			self.feedback_output.setPlainText(f"API 요청 실패: {e}")

	def _format_agent_feedback(self, result: dict[str, object]) -> str:
		name = str(result.get("name") or "")
		char_count = result.get("charCount", 0)
		line_count = result.get("lineCount", 0)
		weak_points = result.get("weakPoints", [])
		suggestions = result.get("suggestions", [])

		lines = [
			f"파일명: {name}",
			f"문서 길이: {char_count}자 / {line_count}줄",
			"",
			"[주요 피드백]",
		]
		if isinstance(weak_points, list):
			for index, point in enumerate(weak_points, start=1):
				lines.append(f"{index}. {point}")
		lines.extend(["", "[개선 제안]"])
		if isinstance(suggestions, list):
			for index, suggestion in enumerate(suggestions, start=1):
				lines.append(f"{index}. {suggestion}")
		return "\n".join(lines).strip()

	def _refresh_stats(self) -> None:
		total = len(self._uploaded_files)
		status = "완료" if total > 0 else "대기"
		self._stats.setText(f"업로드 파일: {total}    |    분석 상태: {status}")

	def _build_feedback(self, file_path: Path) -> str:
		content = self._safe_read(file_path)
		if content.startswith("[읽기 실패]"):
			return f"파일명: {file_path.name}\n\n{content}"

		char_count = len(content)
		line_count = len(content.splitlines()) if content else 0

		weak_points: list[str] = []
		if char_count < 280:
			weak_points.append("문서 길이가 짧아 핵심 근거가 충분하지 않을 수 있습니다.")
		if "출처" not in content and "source" not in content.lower():
			weak_points.append("출처 표기가 보이지 않아 신뢰도 검증이 어렵습니다.")
		if "TODO" in content or "추후" in content:
			weak_points.append("미완료 표기가 포함되어 최종본 품질이 낮아질 수 있습니다.")
		if not weak_points:
			weak_points.append("치명적 문제는 감지되지 않았습니다. 문장 간 연결성과 근거 명확성만 점검하세요.")

		suggestions = [
			"핵심 주장마다 근거 문장 또는 출처를 한 줄 이상 추가하세요.",
			"모호한 표현(예: 일부, 대체로, 빠르게)을 수치/기준으로 구체화하세요.",
			"결론 단락에 실행 항목 2~3개를 명시해 의사결정 가능성을 높이세요.",
		]

		feedback_lines = [
			f"파일명: {file_path.name}",
			f"문서 길이: {char_count}자 / {line_count}줄",
			"",
			"[주요 피드백]",
		]
		for idx, point in enumerate(weak_points, start=1):
			feedback_lines.append(f"{idx}. {point}")

		feedback_lines.append("")
		feedback_lines.append("[개선 제안]")
		for idx, suggestion in enumerate(suggestions, start=1):
			feedback_lines.append(f"{idx}. {suggestion}")

		return "\n".join(feedback_lines)

	def _safe_read(self, file_path: Path) -> str:
		suffix = file_path.suffix.lower()

		if suffix in {".txt", ".md", ".rst", ".log"}:
			return self._read_text_file(file_path)
		if suffix == ".pdf":
			return self._read_pdf(file_path)
		if suffix == ".docx":
			return self._read_docx(file_path)
		if suffix == ".pptx":
			return self._read_pptx(file_path)
		if suffix == ".ppt":
			return self._read_ppt(file_path)
		if suffix == ".hwp":
			return self._read_hwp(file_path)

		return self._read_text_file(file_path)

	def _read_text_file(self, file_path: Path) -> str:
		try:
			return file_path.read_text(encoding="utf-8")
		except UnicodeDecodeError:
			try:
				return file_path.read_text(encoding="cp949", errors="ignore")
			except OSError:
				return "[읽기 실패] 텍스트 파일을 읽을 수 없습니다."
		except OSError:
			return "[읽기 실패] 파일을 읽을 수 없습니다."

	def _read_pdf(self, file_path: Path) -> str:
		if PdfReader is None:
			return "[읽기 실패] PDF 파서를 찾을 수 없습니다. requirements 설치 후 다시 시도하세요."

		try:
			reader = PdfReader(str(file_path))
			extracted = [page.extract_text() or "" for page in reader.pages]
			text = "\n".join(extracted).strip()
			return text or "[읽기 실패] PDF에서 추출 가능한 텍스트가 없습니다."
		except Exception:
			return "[읽기 실패] PDF 파일을 파싱할 수 없습니다."

	def _read_docx(self, file_path: Path) -> str:
		if Document is None:
			return "[읽기 실패] DOCX 파서를 찾을 수 없습니다. requirements 설치 후 다시 시도하세요."

		try:
			doc = Document(str(file_path))
			parts = [p.text for p in doc.paragraphs if p.text.strip()]
			text = "\n".join(parts).strip()
			return text or "[읽기 실패] DOCX에서 추출 가능한 텍스트가 없습니다."
		except Exception:
			return "[읽기 실패] DOCX 파일을 파싱할 수 없습니다."

	def _read_pptx(self, file_path: Path) -> str:
		if Presentation is None:
			return "[읽기 실패] PPTX 파서를 찾을 수 없습니다. requirements 설치 후 다시 시도하세요."

		try:
			prs = Presentation(str(file_path))
			chunks: list[str] = []
			for slide in prs.slides:
				for shape in slide.shapes:
					if hasattr(shape, "text") and shape.text:
						chunks.append(shape.text)
			text = "\n".join(chunks).strip()
			return text or "[읽기 실패] PPTX에서 추출 가능한 텍스트가 없습니다."
		except Exception:
			return "[읽기 실패] PPTX 파일을 파싱할 수 없습니다."

	def _read_ppt(self, file_path: Path) -> str:
		try:
			data = file_path.read_bytes()
		except OSError:
			return "[읽기 실패] PPT 파일을 읽을 수 없습니다."

		utf16_text = data.decode("utf-16le", errors="ignore")
		ansi_text = data.decode("cp949", errors="ignore")
		mixed = f"{utf16_text}\n{ansi_text}"

		candidates = re.findall(r"[A-Za-z0-9가-힣][A-Za-z0-9가-힣\s\-_,.:/()]{5,}", mixed)
		cleaned = [" ".join(c.split()) for c in candidates]
		unique: list[str] = []
		seen: set[str] = set()
		for line in cleaned:
			if line in seen:
				continue
			seen.add(line)
			unique.append(line)

		text = "\n".join(unique[:120]).strip()
		return text or "[읽기 실패] PPT에서 추출 가능한 텍스트를 찾지 못했습니다."

	def _read_hwp(self, file_path: Path) -> str:
		if olefile is None:
			return "[읽기 실패] HWP 파서를 찾을 수 없습니다. requirements 설치 후 다시 시도하세요."

		try:
			if not olefile.isOleFile(str(file_path)):
				return self._read_hwpx_like_zip(file_path)

			ole = olefile.OleFileIO(str(file_path))
			try:
				stream_name = "PrvText"
				if not ole.exists(stream_name):
					return "[읽기 실패] HWP 텍스트 스트림(PrvText)을 찾지 못했습니다."

				with ole.openstream(stream_name) as stream:
					raw = stream.read()

				text = raw.decode("utf-16le", errors="ignore").replace("\x00", "").strip()
				return text or "[읽기 실패] HWP에서 추출 가능한 텍스트가 없습니다."
			finally:
				ole.close()
		except Exception:
			return "[읽기 실패] HWP 파일을 파싱할 수 없습니다."

	def _read_hwpx_like_zip(self, file_path: Path) -> str:
		try:
			with zipfile.ZipFile(file_path) as zf:
				xml_names = [name for name in zf.namelist() if name.lower().endswith(".xml")]
				texts: list[str] = []
				for name in xml_names:
					data = zf.read(name)
					try:
						root = ElementTree.fromstring(data)
					except ElementTree.ParseError:
						continue

					for node in root.iter():
						if node.text and node.text.strip():
							texts.append(node.text.strip())

				merged = "\n".join(texts).strip()
				return merged or "[읽기 실패] HWP/HWPX에서 추출 가능한 텍스트가 없습니다."
		except Exception:
			return "[읽기 실패] HWP/HWPX 파일을 파싱할 수 없습니다."
