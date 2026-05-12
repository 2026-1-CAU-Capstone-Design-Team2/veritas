from __future__ import annotations

from io import BytesIO
import re
from typing import Any
import zipfile
from xml.etree import ElementTree

from fastapi import HTTPException, UploadFile

from ..api_common import new_id
from ..repositories import state_repository as repo
from .agent_runtime import get_runtime

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None  # type: ignore[assignment]

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]

try:
    import olefile
except ImportError:  # pragma: no cover
    olefile = None  # type: ignore[assignment]


async def upload_feedback_files(files: list[UploadFile]) -> dict[str, list[dict[str, str]]]:
    items: list[dict[str, str]] = []
    for file in files:
        file_id = new_id("file")
        file_name = file.filename or "unknown"
        content_type = file.content_type or "application/octet-stream"
        raw = await file.read()
        text = _extract_text(file_name, raw)
        repo.save_feedback_file(file_id, file_name, content_type, text)
        items.append({"fileId": file_id, "name": file_name})
    return {"items": items}


def analyze_feedback(file_ids: list[str]) -> dict[str, str]:
    missing_file_ids = [file_id for file_id in file_ids if repo.get_feedback_file(file_id) is None]
    if missing_file_ids:
        raise HTTPException(status_code=404, detail=f"file(s) not found: {', '.join(missing_file_ids)}")

    analysis_id = new_id("an")
    for file_id in file_ids:
        file_info = repo.get_feedback_file(file_id) or {}
        text = str(file_info.get("text") or "")
        prompt = (
            "다음 피드백 문서의 핵심 약점과 개선 제안을 분석해 주세요. "
            "간결하게 정리하되, 원문에 없는 사실은 만들지 마세요.\n\n"
            f"{text}"
        )
        analysis = get_runtime().answer_chat(prompt, mode="research")
        repo.save_feedback_result(
            file_id,
            {
                "fileId": file_id,
                "name": file_info.get("name", ""),
                "charCount": len(text),
                "lineCount": len(text.splitlines()),
                "weakPoints": [analysis],
                "suggestions": [analysis],
            },
        )

    repo.save_feedback_session(analysis_id, file_ids, "completed")
    return {"analysisId": analysis_id, "status": "completed"}


def get_feedback_result(file_id: str) -> dict[str, Any]:
    file_info = repo.get_feedback_file(file_id)
    if file_info is None:
        raise HTTPException(status_code=404, detail=f"file '{file_id}' not found")

    result = repo.get_feedback_result(file_id)
    if result is None:
        raise HTTPException(status_code=409, detail=f"file '{file_id}' has not been analyzed")
    return result


def clear_feedback_session(session_id: str) -> None:
    repo.clear_feedback_session(session_id)


def _extract_text(file_name: str, raw: bytes) -> str:
    suffix = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    if not raw:
        return ""
    if suffix in {"txt", "md", "rst", "log"}:
        return _decode_text(raw)
    if suffix == "pdf":
        return _extract_pdf(raw)
    if suffix == "docx":
        return _extract_docx(raw)
    if suffix == "pptx":
        return _extract_pptx(raw)
    if suffix == "ppt":
        return _extract_ppt(raw)
    if suffix in {"hwp", "hwpx"}:
        return _extract_hwp(raw)
    return _decode_text(raw)


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8", "cp949", "utf-16le"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _extract_pdf(raw: bytes) -> str:
    if PdfReader is None:
        return _decode_text(raw)
    try:
        reader = PdfReader(BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip()
    except Exception:
        return _decode_text(raw)


def _extract_docx(raw: bytes) -> str:
    if Document is None:
        return _decode_text(raw)
    try:
        doc = Document(BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()
    except Exception:
        return _decode_text(raw)


def _extract_pptx(raw: bytes) -> str:
    if Presentation is None:
        return _decode_text(raw)
    try:
        prs = Presentation(BytesIO(raw))
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks).strip()
    except Exception:
        return _decode_text(raw)


def _extract_ppt(raw: bytes) -> str:
    mixed = f"{raw.decode('utf-16le', errors='ignore')}\n{raw.decode('cp949', errors='ignore')}"
    candidates = re.findall(r"[A-Za-z0-9가-힣][A-Za-z0-9가-힣\s\-_,.:/()]{5,}", mixed)
    unique: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        line = " ".join(candidate.split())
        if line and line not in seen:
            seen.add(line)
            unique.append(line)
    return "\n".join(unique[:120]).strip()


def _extract_hwp(raw: bytes) -> str:
    if olefile is not None:
        try:
            with BytesIO(raw) as buffer:
                if olefile.isOleFile(buffer):
                    buffer.seek(0)
                    ole = olefile.OleFileIO(buffer)
                    try:
                        if ole.exists("PrvText"):
                            with ole.openstream("PrvText") as stream:
                                return stream.read().decode("utf-16le", errors="ignore").replace("\x00", "").strip()
                    finally:
                        ole.close()
        except Exception:
            pass
    return _extract_zip_xml(raw) or _decode_text(raw)


def _extract_zip_xml(raw: bytes) -> str:
    try:
        with zipfile.ZipFile(BytesIO(raw)) as zf:
            texts: list[str] = []
            for name in zf.namelist():
                if not name.lower().endswith(".xml"):
                    continue
                try:
                    root = ElementTree.fromstring(zf.read(name))
                except ElementTree.ParseError:
                    continue
                for node in root.iter():
                    if node.text and node.text.strip():
                        texts.append(node.text.strip())
            return "\n".join(texts).strip()
    except Exception:
        return ""
